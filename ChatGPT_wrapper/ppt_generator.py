from pptx import Presentation
from pptx.util import Inches

def ppt_generator(question, answer):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    # title for question
    title.text = "ChatGPT is answering for your question:"
    subtitle.text = question

    # text box for answer
    blank_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    tf.text = answer
    prs.save('test.pptx')
